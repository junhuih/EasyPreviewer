#!/usr/bin/env python3
from pathlib import Path
from datetime import date
import json
import csv

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.drawing.image import Image as XLImage
from openpyxl.chart import BarChart, Reference
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parent.parent
DEMO_DIR = ROOT / "frontend" / "public" / "demo" / "files"
DEMO_DIR.mkdir(parents=True, exist_ok=True)


def make_brand_image() -> Path:
    image_path = DEMO_DIR / "preview-poster.png"
    width, height = 1280, 720
    image = Image.new("RGB", (width, height), "#f4efe4")
    draw = ImageDraw.Draw(image)

    for idx, color in enumerate(["#0d223f", "#d27b32", "#8bb8e8", "#305d52"]):
        draw.ellipse(
            (80 + idx * 220, 80 + idx * 35, 420 + idx * 220, 420 + idx * 35),
            fill=color,
        )

    draw.rectangle((70, 520, 1210, 620), fill="#0d223f")
    draw.rectangle((70, 620, 1210, 655), fill="#d27b32")

    try:
        title_font = ImageFont.truetype("Helvetica.ttc", 64)
        body_font = ImageFont.truetype("Helvetica.ttc", 28)
    except Exception:
        title_font = ImageFont.load_default()
        body_font = ImageFont.load_default()

    draw.text((90, 520), "Preview-Only Demo Poster", fill="white", font=title_font)
    draw.text(
        (92, 605),
        "Used inside DOCX, XLSX, PPTX, PNG, and JPG sample files.",
        fill="#f7efe0",
        font=body_font,
    )

    image.save(image_path)
    image.convert("RGB").save(DEMO_DIR / "preview-poster.jpg", quality=92)
    return image_path


def make_docx(image_path: Path) -> None:
    doc = Document()
    title = doc.add_heading("EasyPreviewer", level=0)
    title.alignment = 1

    intro = doc.add_paragraph()
    intro.add_run("This DOCX sample validates LibreOffice conversion with ").bold = True
    intro.add_run("headings, styled text, a table, and embedded images.")

    bullet_points = [
        "Bilingual documentation target: zh / en",
        "Preview-only policy: no edit, no export, no upload",
        "Brand target: EasyPreviewer",
    ]
    for item in bullet_points:
        doc.add_paragraph(item, style="List Bullet")

    run = doc.add_paragraph().add_run("Formatted emphasis block")
    run.italic = True
    run.font.color.rgb = RGBColor(0x0D, 0x22, 0x3F)
    run.font.size = Pt(18)

    table = doc.add_table(rows=1, cols=3)
    headers = ["Capability", "Status", "Notes"]
    for idx, text in enumerate(headers):
        table.rows[0].cells[idx].text = text
    rows = [
        ("Office preview", "Supported", "LibreOffice + JODConverter"),
        ("CAD preview", "Unsupported", "Excluded in v1"),
        ("Media transcoding", "Unsupported", "Excluded in v1"),
    ]
    for capability, status, notes in rows:
        cells = table.add_row().cells
        cells[0].text = capability
        cells[1].text = status
        cells[2].text = notes

    doc.add_picture(str(image_path), width=Inches(5.8))
    doc.add_paragraph(f"Generated on {date.today().isoformat()}.")
    doc.save(DEMO_DIR / "sample-complex.docx")


def make_xlsx(image_path: Path) -> None:
    workbook = Workbook()
    issues = workbook.active
    issues.title = "Issue Log"

    header_fill = PatternFill("solid", fgColor="0D223F")
    header_font = Font(color="FFFFFF", bold=True, size=12)
    highlight_fill = PatternFill("solid", fgColor="F2D6BC")
    thin = Side(style="thin", color="D8D8D8")

    issues.merge_cells("A1:H1")
    issues["A1"] = "EasyPreviewer Inspection Workbook"
    issues["A1"].font = Font(size=18, bold=True, color="0D223F")
    issues["A1"].alignment = Alignment(horizontal="center")

    issues["A2"] = "Part"
    issues["B2"] = "Entry Date"
    issues["C2"] = "Trial Date"
    issues["D2"] = "Tool Status"
    issues["E2"] = "Issue"
    issues["F2"] = "Spec"
    issues["G2"] = "Cavity"
    issues["H2"] = "Images"

    for cell in issues[2]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")
        cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)

    data_rows = [
        ["Housing-A", "2026/02/03", "2026/02/05", "OK", "Rough Surface", "Visual", 1],
        ["Housing-B", "2026/02/03", "2026/02/05", "OK", "Flow Line", "Visual", 2],
        ["Housing-C", "2026/02/03", "2026/02/05", "Hold", "Deformation", "Dimensional", 1],
    ]
    for row in data_rows:
        issues.append(row)

    for row in issues.iter_rows(min_row=3, max_row=5, min_col=1, max_col=8):
        for cell in row:
            cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)
            cell.alignment = Alignment(horizontal="center", vertical="center")
        row[4].fill = highlight_fill
        row[7].alignment = Alignment(horizontal="center", vertical="center")

    issues.row_dimensions[3].height = 150
    issues.row_dimensions[4].height = 150
    issues.row_dimensions[5].height = 150

    chart = BarChart()
    chart.title = "Issue Count by Part"
    chart.y_axis.title = "Count"
    chart.x_axis.title = "Part"
    summary_data = [
        ("Housing-A", 1),
        ("Housing-B", 1),
        ("Housing-C", 1),
    ]

    summary = workbook.create_sheet("Summary")
    summary["A1"] = "EasyPreviewer Summary"
    summary["A1"].font = Font(size=18, bold=True, color="0D223F")
    summary.append(["Part", "Issue Count"])
    for row in summary_data:
        summary.append(list(row))
    for cell in summary[2]:
        cell.fill = header_fill
        cell.font = header_font

    data = Reference(summary, min_col=2, max_col=2, min_row=2, max_row=5)
    cats = Reference(summary, min_col=1, min_row=3, max_row=5)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 6
    chart.width = 11
    summary.add_chart(chart, "D2")

    issues.add_image(XLImage(str(image_path)), "H3")

    for column, width in {"A": 18, "B": 14, "C": 14, "D": 14, "E": 18, "F": 14, "G": 10, "H": 32}.items():
        issues.column_dimensions[column].width = width

    gallery = workbook.create_sheet("Photo Board")
    gallery["A1"] = "EasyPreviewer Photo Board"
    gallery["A1"].font = Font(size=18, bold=True, color="0D223F")
    gallery["A3"] = "Reference image"
    gallery.add_image(XLImage(str(image_path)), "A5")

    for cell in summary[2]:
        cell.fill = header_fill
        cell.font = header_font

    workbook.save(DEMO_DIR / "sample-complex.xlsx")


def make_pptx(image_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "EasyPreviewer Demo Deck"
    title_slide.placeholders[1].text = "Complex PPTX sample with layout, images, tables, and bilingual content."

    agenda = prs.slides.add_slide(prs.slide_layouts[5])
    agenda.shapes.title.text = "Agenda"
    tx = agenda.shapes.add_textbox(PptInches(0.8), PptInches(1.4), PptInches(5.2), PptInches(3.8))
    tf = tx.text_frame
    tf.word_wrap = True
    for idx, item in enumerate([
        "1. Product rules",
        "2. Supported file types",
        "3. Unsupported-by-policy list",
        "4. Bilingual rollout",
    ]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = PptPt(24)
        p.font.bold = idx == 0

    agenda.shapes.add_picture(str(image_path), PptInches(7.1), PptInches(1.2), width=PptInches(5.2))

    matrix = prs.slides.add_slide(prs.slide_layouts[5])
    matrix.shapes.title.text = "Support Matrix"
    table = matrix.shapes.add_table(4, 3, PptInches(0.7), PptInches(1.5), PptInches(8.8), PptInches(2.5)).table
    table.columns[0].width = PptInches(2.7)
    table.columns[1].width = PptInches(1.8)
    table.columns[2].width = PptInches(4.3)
    values = [
        ("Capability", "Status", "Notes"),
        ("Office", "Supported", "LibreOffice conversion"),
        ("PDF / Text / Images", "Supported", "Read-only rendering"),
        ("CAD / Media / TIFF", "Unsupported", "Removed by policy"),
    ]
    for row_idx, row_values in enumerate(values):
        for col_idx, value in enumerate(row_values):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = PptPt(18 if row_idx == 0 else 16)
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = PptRGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptRGBColor(13, 34, 63) if row_idx == 0 else PptRGBColor(244, 239, 228)

    bilingual = prs.slides.add_slide(prs.slide_layouts[1])
    bilingual.shapes.title.text = "Bilingual UX"
    body = bilingual.placeholders[1].text_frame
    body.clear()
    p1 = body.paragraphs[0]
    p1.text = "English: Preview only means no editing, no export, and no save-back."
    p1.font.size = PptPt(24)
    p1.alignment = PP_ALIGN.LEFT
    p2 = body.add_paragraph()
    p2.text = "中文：仅预览意味着不提供编辑、导出与保存回源文件。"
    p2.font.size = PptPt(24)
    p2.font.color.rgb = PptRGBColor(48, 93, 82)

    prs.save(DEMO_DIR / "sample-complex.pptx")


def make_textual_assets() -> None:
    (DEMO_DIR / "sample-readme.md").write_text(
        "# Demo Markdown\n\n"
        "This markdown sample is used by the preview shell.\n\n"
        "- English: preview only\n"
        "- 中文：仅支持预览\n\n"
        "```ts\n"
        "export const previewMode = 'read-only'\n"
        "```\n",
        encoding="utf-8",
    )
    (DEMO_DIR / "sample-notes.txt").write_text(
        "Preview-only checklist\n"
        "1. No edit\n"
        "2. No export\n"
        "3. No upload\n"
        "4. zh / en only\n",
        encoding="utf-8",
    )
    (DEMO_DIR / "sample-code.ts").write_text(
        "type PreviewStatus = 'READY' | 'PROCESSING' | 'UNSUPPORTED' | 'FAILED'\n"
        "export const supportedLocales = ['zh', 'en'] as const\n",
        encoding="utf-8",
    )
    (DEMO_DIR / "sample-data.json").write_text(
        json.dumps(
            {
                "product": "modern-preview-viewer",
                "brand": "EasyPreviewer",
                "previewOnly": True,
                "locales": ["zh", "en"],
                "unsupported": ["cad", "tiff-conversion", "media-transcoding"],
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    with (DEMO_DIR / "sample-grid.csv").open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["team", "capability", "status", "owner"])
        writer.writerow(["platform", "office-preview", "supported", "backend"])
        writer.writerow(["platform", "cad-preview", "unsupported", "policy"])
        writer.writerow(["frontend", "markdown-preview", "supported", "ui"])

    (DEMO_DIR / "sample-badge.svg").write_text(
        """
<svg xmlns="http://www.w3.org/2000/svg" width="720" height="220" viewBox="0 0 720 220">
  <rect width="720" height="220" rx="24" fill="#0d223f"/>
  <circle cx="108" cy="110" r="54" fill="#d27b32"/>
  <text x="190" y="96" fill="#f8f1e7" font-size="42" font-family="Georgia">Preview-Only Demo</text>
  <text x="192" y="146" fill="#bad8f6" font-size="24" font-family="Helvetica">zh / en · Apache-2.0 · LibreOffice</text>
</svg>
""".strip()
        + "\n",
        encoding="utf-8",
    )


def main() -> None:
    image_path = make_brand_image()
    make_docx(image_path)
    make_xlsx(image_path)
    make_pptx(image_path)
    make_textual_assets()


if __name__ == "__main__":
    main()
